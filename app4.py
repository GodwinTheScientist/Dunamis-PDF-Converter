import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF
import re
from io import BytesIO

logo_url = "https://i.postimg.cc/sxSLVk2D/church-logo-cmyk-1-white.png"
bg_image_url = "https://images.unsplash.com/photo-1530688957198-8570b1819eeb?q=80&w=2114&auto=format&fit=crop&ixlib=rb-4.1.0&ixid=M3wxMjA3fDB8MHxwaG90by1wYWdlfHx8fGVufDB8fHx8fA%3D%3D"

st.markdown(f"""
    <style>
    .stApp {{
        background: radial-gradient(circle, rgba(15, 23, 42, 0.9) 0%, rgba(0, 0, 20, 0.98) 100%), 
                    url("{bg_image_url}");
        background-size: cover;
        background-position: center;
        background-attachment: fixed;
    }}
    .header-box {{
        text-align: center;
        padding: 90px 0 40px 0;
        min-height: 180px;
        display: flex;
        flex-direction: column;
        align-items: center;
        justify-content: center;
    }}
    .mini-logo {{
        width: 120px;
        filter: drop-shadow(0px 0px 12px rgba(255, 255, 255, 0.2));
        margin-bottom: 15px;
    }}
    .main-title {{
        color: #FFD700 !important;
        font-size: 2.4rem !important;
        font-weight: 900 !important;
        margin: 0 !important;
        text-shadow: 0 4px 12px rgba(0,0,0,0.7);
    }}
    .subtitle {{
        color: #94A3B8;
        font-size: 1.1rem;
        margin: 8px 0 0 0;
    }}
    [data-testid="stTabList"] {{
        display: flex !important;
        justify-content: center !important;
        max-width: 600px !important;
        margin: 0 auto 20px auto !important;
        background: rgba(255,255,255,0.06) !important;
        border-radius: 12px !important;
        padding: 6px !important;
        backdrop-filter: blur(12px) !important;
    }}
    [data-testid="stTab"] {{
        color: #CBD5E1 !important;
        font-weight: 600;
        padding: 10px 24px !important;
        font-size: 1rem !important;
    }}
    [aria-selected="true"] {{
        background: rgba(255,215,0,0.25) !important;
        color: #FFD700 !important;
        border-radius: 8px !important;
        padding: 6px !important;
    }}
    .block-container {{
        max-width: 900px !important;
        margin: 0 auto !important;
        padding: 0 20px !important;
    }}
    .stTabs > div > div:has(> *) {{
        background: rgba(255,255,255,0.08) !important;
        backdrop-filter: blur(16px) !important;
        border-radius: 16px !important;
        border: 1px solid rgba(255,255,255,0.12) !important;
        box-shadow: 0 8px 24px rgba(0,0,0,0.5) !important;
        padding: 25px 20px !important;
        margin: 15px auto !important;
    }}
    .generate-container {{
        text-align: center;
        margin: 40px 0 40px 0;
    }}
    .metric-card {{
        background: rgba(255,255,255,0.08) !important;
        backdrop-filter: blur(16px) !important;
        border-radius: 16px !important;
        border: 1px solid rgba(255,255,255,0.12) !important;
        box-shadow: 0 8px 20px rgba(0,0,0,0.5) !important;
        padding: 20px !important;
        text-align: center;
    }}
    .metric-card h3 {{
        color: #FFD700 !important;
        font-size: 2.4rem !important;
        margin: 0 !important;
    }}
    .metric-card p {{
        color: #CBD5E1 !important;
        font-size: 0.95rem !important;
        margin: 6px 0 0 !important;
    }}
    </style>

    <div class="header-box">
        <img src="{logo_url}" class="mini-logo">
        <h1 class="main-title">Dunamis Prayer Converter</h1>
        <p class="subtitle">PDF to PPTX Dashboard</p>
    </div>
""", unsafe_allow_html=True)

# ── Metrics Syncing ─────────────────────────────────────────────────────────
if 'total_prayers_count' not in st.session_state:
    st.session_state.total_prayers_count = "-"
if 'total_sessions_count' not in st.session_state:
    st.session_state.total_sessions_count = "-"

cols = st.columns(3)
with cols[0]:
    num_pdfs = len(st.session_state.get('uploaded_files', []))
    st.markdown(f"<div class='metric-card'><h3>{num_pdfs}</h3><p>Total PDFs</p></div>", unsafe_allow_html=True)
with cols[1]:
    st.markdown(f"<div class='metric-card'><h3>{st.session_state.total_prayers_count}</h3><p>Prayers</p></div>", unsafe_allow_html=True)
with cols[2]:
    st.markdown(f"<div class='metric-card'><h3>{st.session_state.total_sessions_count}</h3><p>Sessions</p></div>", unsafe_allow_html=True)

tab1, tab2 = st.tabs(["📁 Upload", "🎨 Customise"])

with tab1:
    uploaded_files = st.file_uploader("Upload PDFs", type=["pdf"], accept_multiple_files=True)
    if uploaded_files:
        st.session_state.uploaded_files = uploaded_files
        st.success(f"Uploaded {len(uploaded_files)} PDF(s)")
        
        prayers_found = 0
        for f in uploaded_files:
            try:
                doc = fitz.open(stream=f.getvalue(), filetype="pdf")
                full_text = "".join(p.get_text("text") for p in doc)
                prayers_found += len(re.findall(r"^\s*(?:Prayer Point\s*\d+|\d+\.)", full_text, re.M | re.I))
            except:
                pass
        st.session_state.total_prayers_count = prayers_found if prayers_found > 0 else "Detected"
        st.session_state.total_sessions_count = len(uploaded_files)

with tab2:
    col_left, col_right = st.columns([1, 1])
    with col_left:
        bg_option = st.radio("Background", ["Dark Navy", "Black", "Deep Purple", "Custom"])
        if bg_option == "Dark Navy":
            bg_rgb = (10, 20, 60)
        elif bg_option == "Black":
            bg_rgb = (0, 0, 0)
        elif bg_option == "Deep Purple":
            bg_rgb = (25, 0, 50)
        else:
            bg_hex = st.color_picker("Custom RGB", "#0A143C")
            bg_rgb = tuple(int(bg_hex[i:i+2], 16) for i in (1, 3, 5))

        header_color = st.color_picker("Header", "#FFD700")
        body_color = st.color_picker("Body", "#FFFFFF")

    with col_right:
        header_size = st.slider("Header size", 40, 100, 68)
        body_size = st.slider("Maximum Body size", 40, 100, 60)
        text_case = st.selectbox("Text case", ["Original", "UPPERCASE", "lowercase", "Title Case"])

# ── PowerPoint Generation Engine ────────────────────────────────────────────
st.markdown('<div class="generate-container">', unsafe_allow_html=True)
if st.button("🚀 Generate & Download PPTX", key="generate"):
    if 'uploaded_files' not in st.session_state or not st.session_state.uploaded_files:
        st.error("Upload PDFs first.")
    else:
        with st.spinner("Processing Presentation Slides..."):
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            def set_bg(slide):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_rgb)

            def clean_text_block(txt):
                if not txt:
                    return ""
                txt = re.sub(r"P\s*a\s*g\s*e\s*\d+\s*\|\s*\d+", "", txt, flags=re.I)
                txt = re.sub(r"Page\s*\d+", "", txt, flags=re.I)
                txt = re.sub(r"Dunamis Bible Church.*", "", txt, flags=re.I)
                txt = re.sub(r"\(AKA.* Charity.*", "", txt, flags=re.I)
                return txt.strip()

            def add_cover_centered(slide, l, t, w, h, text, size, color_hex, bold=False):
                tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = text
                run.font.size = Pt(size)
                r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (1, 3, 5))
                run.font.color.rgb = RGBColor(r, g, b)
                run.font.bold = bold

            # ── Unified Layout Engine to Prevent Spill Over Truncations ──
            def add_fluid_prayer_slide(slide, num, body_text, max_b_size, h_size, h_color, b_color):
                # We open one massive box across the entire safe viewport width
                left = Inches(0.8)
                top = Inches(0.6)
                width = Inches(11.733)
                height = Inches(6.3) # Total allowed running height space before spilling

                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.word_wrap = True
                
                # Wipe margins entirely to drop padding calculations
                tf.margin_top = Inches(0)
                tf.margin_bottom = Inches(0)
                tf.margin_left = Inches(0)
                tf.margin_right = Inches(0)

                # Loop to safely fit text within frame boundaries
                current_size = max_b_size
                char_len = len(body_text)

                # Dynamic layout fallback tuning based on string density parameters
                if char_len > 300:
                    current_size = min(current_size, 38)
                elif char_len > 180:
                    current_size = min(current_size, 46)

                while current_size >= 24:
                    tf.clear()
                    
                    # 1. Header Paragraph
                    p1 = tf.paragraphs[0]
                    p1.alignment = PP_ALIGN.CENTER
                    p1.space_after = Pt(24) # Controlled safe gap separating title and blocks
                    
                    run1 = p1.add_run()
                    run1.text = f"Prayer Point {num}"
                    run1.font.size = Pt(h_size)
                    r, g, b = tuple(int(h_color[i:i+2], 16) for i in (1, 3, 5))
                    run1.font.color.rgb = RGBColor(r, g, b)
                    run1.font.bold = True
                    
                    # 2. Body Text Paragraph
                    p2 = tf.add_paragraph()
                    p2.alignment = PP_ALIGN.CENTER
                    p2.line_spacing = 1.15
                    
                    run2 = p2.add_run()
                    run2.text = body_text
                    run2.font.size = Pt(current_size)
                    br, bg, bb = tuple(int(b_color[i:i+2], 16) for i in (1, 3, 5))
                    run2.font.color.rgb = RGBColor(br, bg, bb)
                    run2.font.bold = True
                    
                    # Estimate line usage footprint to actively check for text overflow
                    estimated_lines = (char_len * (current_size * 0.55)) / (11.733 * 72)
                    estimated_height = (h_size + 24 + (estimated_lines * current_size * 1.2)) / 72
                    
                    if estimated_height <= 6.0:
                        break
                    current_size -= 4  # Scale down until it completely fits inside the slide frame

            for idx, file in enumerate(st.session_state.uploaded_files):
                doc = fitz.open(stream=file.getvalue(), filetype="pdf")
                text = "".join(page.get_text("text") for page in doc)
                lines = [l.strip() for l in text.split("\n") if l.strip()]

                title = file.name.replace(".pdf", "").replace("_", " ")
                if "FRIDAY" in text.upper():
                    title = "Friday Prayer Session"
                elif "SATURDAY" in text.upper():
                    title = "Saturday Prayer Session"

                if idx > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    set_bg(slide)
                    add_cover_centered(slide, 0.8, 2.8, 11.7, 2.0, f"--- Next Session ---\n{title}", 56, header_color, True)

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                set_bg(slide)
                add_cover_centered(slide, 0.8, 1.2, 11.7, 1.8, f"Dunamis Bible Church", 56, header_color, True)
                add_cover_centered(slide, 1.2, 4.8, 11.0, 1.2, "PRAYER POINTS", 50, "#CCCCCC")

                prayers = []
                current = ""
                
                for line in lines:
                    if any(m in line.lower() for m in ["charity no", "dunamis centre", "northmoor", "manchester m12", "info@", "+44", "prayer session"]):
                        continue
                    if re.match(r"^P\s*a\s*g\s*e\s*\d+", line, re.I) or line.strip() == "Dunamis Bible Church" or "PRAYER POINTS" in line:
                        continue
                    if any(x in line for x in ["IJN=", "ITNJ=", "ITMNJ=", "ITNJCN=", "(KJV)"]):
                        continue
                    if not current and not re.match(r"^\d+\.", line) and not re.match(r"^Prayer Point\s*\d+", line, re.I):
                        continue
                        
                    if re.match(r"^\d+\.", line) or re.match(r"^Prayer Point\s*\d+", line, re.I):
                        if current:
                            prayers.append(current.strip())
                        current = line
                    elif current:
                        current += " " + line
                if current:
                    prayers.append(current.strip())

                for prayer in prayers:
                    m = re.match(r"^(?:Prayer Point\s*)?(\d+)[\.\s]*(.*)", prayer, re.DOTALL | re.I)
                    if m:
                        num, text_content = m.groups()
                        text_content = clean_text_block(text_content)
                        if not text_content:
                            continue

                        if text_case == "UPPERCASE":
                            text_content = text_content.upper()
                        elif text_case == "lowercase":
                            text_content = text_content.lower()
                        elif text_case == "Title Case":
                            text_content = text_content.title()

                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        set_bg(slide)
                        
                        # Process using our layout engine to prevent spilling
                        add_fluid_prayer_slide(slide, num, text_content, body_size, header_size, header_color, body_color)

            bio = BytesIO()
            prs.save(bio)
            bio.seek(0)

            st.success("✅ Layout Generated Successfully!")
            st.download_button(
                label="⬇ Download PPTX",
                data=bio,
                file_name="Dunamis_Prayer_Points.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

st.markdown('</div>', unsafe_allow_html=True)